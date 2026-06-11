#!/usr/bin/env python3
"""Convert Reveal.js HTML presentations to PowerPoint (.pptx).

Usage:
    python3 convert_to_pptx.py <input.html> [output.pptx]

If output is omitted, saves alongside the HTML file with same name.
"""

import sys
import re
from pathlib import Path
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Theme colors (from Onyx Armor CSS vars) ───
BG_DARK = RGBColor(0x0A, 0x0A, 0x0A)
BG_CARD = RGBColor(0x11, 0x11, 0x11)
ACCENT = RGBColor(0x7B, 0x8F, 0xA8)
ACCENT_LIGHT = RGBColor(0x9B, 0xB0, 0xC9)
SUCCESS = RGBColor(0x7B, 0x9A, 0x7B)
WARNING = RGBColor(0xC9, 0xA8, 0x5B)
ERROR = RGBColor(0xC9, 0x7B, 0x7B)
TEXT_PRIMARY = RGBColor(0xF5, 0xF4, 0xF0)
TEXT_SECONDARY = RGBColor(0xA8, 0xA8, 0xA8)
TEXT_MUTED = RGBColor(0x66, 0x66, 0x66)
BORDER = RGBColor(0x2A, 0x2A, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


class HTMLTextExtractor(HTMLParser):
    """Extract plain text from HTML fragments."""
    def __init__(self):
        super().__init__()
        self._text = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in ('br',):
            self._text.append('\n')
        if tag == 'style':
            self._skip = True

    def handle_endtag(self, tag):
        if tag == 'style':
            self._skip = False

    def handle_data(self, data):
        if not self._skip:
            self._text.append(data)

    def get_text(self):
        return ''.join(self._text).strip()


def html_to_text(html_str):
    """Strip HTML tags, convert <br> to newlines."""
    extractor = HTMLTextExtractor()
    extractor.feed(html_str)
    return extractor.get_text()


def extract_css_vars(html):
    """Extract CSS custom properties from the HTML."""
    colors = {}
    var_pattern = re.compile(r'--([\w-]+):\s*([^;]+);')
    style_match = re.search(r'<style>(.*?)</style>', html, re.DOTALL)
    if style_match:
        for m in var_pattern.finditer(style_match.group(1)):
            colors[m.group(1)] = m.group(2).strip()
    return colors


def parse_slides(html):
    """Parse <section> blocks into slide data structures."""
    sections = re.findall(r'<section>(.*?)</section>', html, re.DOTALL)
    slides = []

    for section in sections:
        slide = {'raw': section, 'elements': []}

        # Detect slide type
        if 'slide--section' in section and 'slide--center' in section:
            slide['type'] = 'section'
        elif 'split' in section and '<table' in section:
            slide['type'] = 'split-table'
        elif 'split' in section:
            slide['type'] = 'split'
        elif 'grid-3' in section:
            slide['type'] = 'grid3'
        elif 'grid-2' in section:
            slide['type'] = 'grid2'
        elif 'compare' in section:
            slide['type'] = 'compare'
        elif 'pipeline' in section:
            slide['type'] = 'pipeline'
        elif 'timeline' in section:
            slide['type'] = 'timeline'
        elif 'metrics' in section:
            slide['type'] = 'metrics'
        elif 'contact-grid' in section:
            slide['type'] = 'section'  # closing CTA
        else:
            slide['type'] = 'content'

        # Extract kicker
        kicker_m = re.search(r'class="kicker"[^>]*>(.*?)</p>', section, re.DOTALL)
        if kicker_m:
            slide['kicker'] = html_to_text(kicker_m.group(1))

        # Extract title (h1 or h2)
        title_m = re.search(r'<h[12][^>]*>(.*?)</h[12]>', section, re.DOTALL)
        if title_m:
            slide['title'] = html_to_text(title_m.group(1))

        # Extract lead text
        lead_m = re.search(r'class="lead"[^>]*>(.*?)</p>', section, re.DOTALL)
        if lead_m:
            slide['lead'] = html_to_text(lead_m.group(1))

        # Extract metrics
        if slide['type'] == 'metrics':
            metrics = re.findall(
                r'class="metric__value">(.*?)</div>\s*<div class="metric__label">(.*?)</div>',
                section, re.DOTALL
            )
            slide['metrics'] = [(html_to_text(v), html_to_text(l)) for v, l in metrics]

        # Extract table
        if '<table' in section:
            rows = []
            thead = re.search(r'<thead>(.*?)</thead>', section, re.DOTALL)
            if thead:
                ths = re.findall(r'<th[^>]*>(.*?)</th>', thead.group(1), re.DOTALL)
                rows.append([html_to_text(th) for th in ths])
            tbody = re.search(r'<tbody>(.*?)</tbody>', section, re.DOTALL)
            if tbody:
                for tr in re.findall(r'<tr>(.*?)</tr>', tbody.group(1), re.DOTALL):
                    tds = re.findall(r'<td[^>]*>(.*?)</td>', tr, re.DOTALL)
                    rows.append([html_to_text(td) for td in tds])
            slide['table'] = rows

        # Extract list items
        lis = re.findall(r'<li[^>]*>(.*?)</li>', section, re.DOTALL)
        if lis:
            items = []
            for li in lis:
                strong_m = re.search(r'<strong>(.*?)</strong>', li, re.DOTALL)
                text = html_to_text(li)
                if strong_m:
                    items.append({
                        'bold': html_to_text(strong_m.group(1)),
                        'text': text
                    })
                else:
                    items.append({'bold': '', 'text': text})
            slide['list_items'] = items

        # Extract cards
        if slide['type'] in ('grid2', 'grid3'):
            cards = []
            for card_m in re.finditer(r'class="card[^"]*">(.*?)</div>\s*</div>', section, re.DOTALL):
                card_html = card_m.group(1)
                card_title = re.search(r'class="card__title">(.*?)</h4>', card_html, re.DOTALL)
                card_text = re.search(r'class="card__text"[^>]*>(.*?)</p>', card_html, re.DOTALL)
                cards.append({
                    'title': html_to_text(card_title.group(1)) if card_title else '',
                    'text': html_to_text(card_text.group(1)) if card_text else ''
                })
            slide['cards'] = cards

        # Extract pipeline steps
        if slide['type'] == 'pipeline':
            steps = []
            for row_m in re.finditer(r'pipeline__source">(.*?)</div>.*?pipeline__step[^"]*">(.*?)</div>', section, re.DOTALL):
                source = html_to_text(row_m.group(1))
                step = html_to_text(row_m.group(2))
                steps.append((source, step))
            slide['pipeline_steps'] = steps

        # Extract timeline items
        if slide['type'] == 'timeline':
            items = []
            for tm in re.finditer(r'timeline__year">(.*?)</div>.*?timeline__content">(.*?)</div>', section, re.DOTALL):
                items.append((html_to_text(tm.group(1)), html_to_text(tm.group(2))))
            slide['timeline_items'] = items

        # Extract compare columns
        if slide['type'] == 'compare':
            cols = []
            for col_m in re.finditer(r'compare__label">(.*?)</div>(.*?)</div>\s*</div>', section, re.DOTALL):
                label = html_to_text(col_m.group(1))
                items_html = col_m.group(2)
                items = [html_to_text(li) for li in re.findall(r'<li[^>]*>(.*?)</li>', items_html, re.DOTALL)]
                cols.append({'label': label, 'items': items})
            slide['compare_cols'] = cols

        # Extract callout/note boxes
        callout_m = re.search(r'border-left:\s*3px solid[^"]*">\s*<p[^>]*>(.*?)</p>', section, re.DOTALL)
        if callout_m:
            slide['callout'] = html_to_text(callout_m.group(1))

        # Extract contact grid
        contacts = re.findall(
            r'contact-item__label">(.*?)</div>\s*<div class="contact-item__value"[^>]*>(.*?)</div>',
            section, re.DOTALL
        )
        if contacts:
            slide['contacts'] = [(html_to_text(l), html_to_text(v)) for l, v in contacts]

        slides.append(slide)

    return slides


def set_slide_bg(slide, color=BG_DARK):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_bg(slide, left, top, width, height, color=BG_CARD):
    """Add a rounded rectangle background shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape


def build_pptx(slides_data, output_path):
    """Build the PowerPoint file from parsed slide data."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    MARGIN = Inches(0.8)
    CONTENT_W = prs.slide_width - 2 * MARGIN
    HALF_W = Inches(5.5)

    for i, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide)

        stype = sd.get('type', 'content')
        kicker = sd.get('kicker', '')
        title = sd.get('title', '')
        lead = sd.get('lead', '')

        # ── Section / Cover slides ──
        if stype == 'section':
            y = Inches(2.0)
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.4),
                            kicker.upper(), font_size=11, color=ACCENT,
                            alignment=PP_ALIGN.CENTER, font_name='Calibri')
                y += Inches(0.5)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(1.5),
                            title, font_size=40, color=TEXT_PRIMARY, bold=True,
                            alignment=PP_ALIGN.CENTER, font_name='Calibri')
                y += Inches(1.5)
            if lead:
                add_textbox(slide, Inches(2.5), y, Inches(8), Inches(1.0),
                            lead, font_size=16, color=TEXT_SECONDARY,
                            alignment=PP_ALIGN.CENTER)
                y += Inches(1.2)

            # Subtitle line (date or other small text)
            subtitle_m = re.search(r'font-size:\s*0\.75rem[^>]*>(.*?)</p>', sd['raw'], re.DOTALL)
            if subtitle_m:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.4),
                            html_to_text(subtitle_m.group(1)), font_size=10,
                            color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

            # Contact grid
            if sd.get('contacts'):
                y = Inches(5.2)
                cx = Inches(3.0)
                for label, value in sd['contacts']:
                    add_shape_bg(slide, cx, y, Inches(2.2), Inches(0.8))
                    add_textbox(slide, cx + Inches(0.15), y + Inches(0.1),
                                Inches(1.9), Inches(0.3), label, font_size=9,
                                color=TEXT_MUTED)
                    add_textbox(slide, cx + Inches(0.15), y + Inches(0.4),
                                Inches(1.9), Inches(0.3), value, font_size=11,
                                color=ACCENT)
                    cx += Inches(2.5)
            continue

        # ── Header bar ──
        header_y = Inches(0.4)
        add_textbox(slide, MARGIN, header_y, Inches(3), Inches(0.35),
                    'Onyx Armor', font_size=12, color=TEXT_SECONDARY,
                    font_name='Georgia')
        # Header line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN, Inches(0.85), CONTENT_W, Pt(1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER
        line.line.fill.background()

        content_top = Inches(1.1)

        # ── Metrics slide ──
        if stype == 'metrics':
            if kicker:
                add_textbox(slide, MARGIN, content_top, HALF_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=ACCENT)
                content_top += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, content_top, HALF_W, Inches(1.0),
                            title, font_size=32, color=TEXT_PRIMARY, bold=True)
                content_top += Inches(1.0)
            if lead:
                add_textbox(slide, MARGIN, content_top, Inches(8), Inches(0.6),
                            lead, font_size=14, color=TEXT_SECONDARY)
                content_top += Inches(0.8)

            if sd.get('metrics'):
                mx = MARGIN
                mw = Inches(2.5)
                for val, label in sd['metrics']:
                    add_shape_bg(slide, mx, content_top, mw, Inches(1.2))
                    add_textbox(slide, mx + Inches(0.2), content_top + Inches(0.15),
                                mw - Inches(0.4), Inches(0.7), val,
                                font_size=36, color=TEXT_PRIMARY, bold=True,
                                alignment=PP_ALIGN.CENTER)
                    add_textbox(slide, mx + Inches(0.2), content_top + Inches(0.8),
                                mw - Inches(0.4), Inches(0.3), label,
                                font_size=11, color=TEXT_MUTED,
                                alignment=PP_ALIGN.CENTER)
                    mx += mw + Inches(0.3)
            continue

        # ── Split layouts ──
        if stype in ('split', 'split-table'):
            left_x = MARGIN
            right_x = Inches(6.8)
            y = content_top

            # Left column: kicker + title + lead
            if kicker:
                add_textbox(slide, left_x, y, HALF_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=ACCENT)
                y += Inches(0.4)
            if title:
                add_textbox(slide, left_x, y, HALF_W, Inches(1.2),
                            title, font_size=28, color=TEXT_PRIMARY, bold=True)
                y += Inches(1.1)
            if lead:
                add_textbox(slide, left_x, y, HALF_W, Inches(1.5),
                            lead, font_size=13, color=TEXT_SECONDARY)
                y += Inches(1.2)

            # Callout box on left
            if sd.get('callout'):
                add_shape_bg(slide, left_x, y, HALF_W, Inches(0.8))
                # Accent bar
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left_x, y, Pt(3), Inches(0.8)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = ACCENT
                bar.line.fill.background()
                add_textbox(slide, left_x + Inches(0.2), y + Inches(0.1),
                            HALF_W - Inches(0.4), Inches(0.6),
                            sd['callout'], font_size=10, color=TEXT_SECONDARY)

            # Right column: table or list
            ry = content_top + Inches(0.2)
            if sd.get('table'):
                table_data = sd['table']
                rows = len(table_data)
                cols = len(table_data[0]) if table_data else 2
                tw = Inches(5.5)
                th = Inches(0.4) * rows
                tbl = slide.shapes.add_table(rows, cols, right_x, ry, tw, th).table

                for r, row in enumerate(table_data):
                    for c, cell_text in enumerate(row):
                        cell = tbl.cell(r, c)
                        cell.text = cell_text
                        cell.fill.solid()
                        if r == 0:
                            cell.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
                        else:
                            cell.fill.fore_color.rgb = BG_CARD
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(10)
                            paragraph.font.color.rgb = TEXT_PRIMARY if r == 0 else TEXT_SECONDARY
                            paragraph.font.bold = (r == 0)
                            paragraph.font.name = 'Calibri'

            elif sd.get('list_items'):
                for item in sd['list_items']:
                    text = item['text']
                    txBox = slide.shapes.add_textbox(right_x, ry, Inches(5.5), Inches(0.8))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]

                    if item['bold']:
                        run_b = p.add_run()
                        run_b.text = item['bold'] + '\n'
                        run_b.font.size = Pt(12)
                        run_b.font.color.rgb = TEXT_PRIMARY
                        run_b.font.bold = True
                        run_b.font.name = 'Calibri'
                        # Rest of text (after bold part)
                        rest = text.replace(item['bold'], '', 1).strip()
                        if rest:
                            run_r = p.add_run()
                            run_r.text = rest
                            run_r.font.size = Pt(10)
                            run_r.font.color.rgb = TEXT_MUTED
                            run_r.font.name = 'Calibri'
                    else:
                        p.text = text
                        p.font.size = Pt(11)
                        p.font.color.rgb = TEXT_SECONDARY
                        p.font.name = 'Calibri'
                    ry += Inches(0.85)
            continue

        # ── Compare slide ──
        if stype == 'compare':
            y = content_top
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=ACCENT)
                y += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                            title, font_size=28, color=TEXT_PRIMARY, bold=True)
                y += Inches(0.9)

            cols = sd.get('compare_cols', [])
            cx = MARGIN
            col_w = Inches(5.5)
            for ci, col in enumerate(cols):
                is_after = ci == 1
                # Column background
                add_shape_bg(slide, cx, y, col_w, Inches(4.0),
                             RGBColor(0x0F, 0x14, 0x0F) if is_after else BG_CARD)
                # Label
                label_color = SUCCESS if is_after else WARNING
                add_textbox(slide, cx + Inches(0.2), y + Inches(0.15),
                            col_w - Inches(0.4), Inches(0.35),
                            col['label'], font_size=13, color=label_color, bold=True)
                # Items
                iy = y + Inches(0.6)
                for item in col['items']:
                    bullet_color = SUCCESS if is_after else TEXT_SECONDARY
                    prefix = '  ' if is_after else '  '
                    add_textbox(slide, cx + Inches(0.3), iy,
                                col_w - Inches(0.6), Inches(0.35),
                                '• ' + item, font_size=11, color=bullet_color)
                    iy += Inches(0.45)
                cx += col_w + Inches(0.5)
            continue

        # ── Pipeline slide ──
        if stype == 'pipeline':
            y = content_top
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=ACCENT)
                y += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                            title, font_size=28, color=TEXT_PRIMARY, bold=True)
                y += Inches(0.9)

            for source, step in sd.get('pipeline_steps', []):
                add_shape_bg(slide, MARGIN, y, CONTENT_W, Inches(0.55))
                add_textbox(slide, MARGIN + Inches(0.2), y + Inches(0.08),
                            Inches(2.5), Inches(0.35), source,
                            font_size=10, color=ACCENT, bold=True)
                add_textbox(slide, Inches(3.8), y + Inches(0.08),
                            Inches(0.4), Inches(0.35), '\u2192',
                            font_size=12, color=TEXT_MUTED,
                            alignment=PP_ALIGN.CENTER)
                add_textbox(slide, Inches(4.3), y + Inches(0.08),
                            Inches(8), Inches(0.35), step,
                            font_size=10, color=TEXT_SECONDARY)
                y += Inches(0.65)
            continue

        # ── Timeline slide ──
        if stype == 'timeline':
            y = content_top
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=ACCENT)
                y += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                            title, font_size=28, color=TEXT_PRIMARY, bold=True)
                y += Inches(1.0)

            for label, content in sd.get('timeline_items', []):
                # Dot
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, Inches(2.2), y + Inches(0.12), Pt(10), Pt(10)
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = ACCENT
                dot.line.fill.background()
                # Vertical line
                if label != sd['timeline_items'][-1][0]:
                    vline = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(2.24), y + Inches(0.35), Pt(2), Inches(0.4)
                    )
                    vline.fill.solid()
                    vline.fill.fore_color.rgb = BORDER
                    vline.line.fill.background()

                add_textbox(slide, MARGIN, y, Inches(1.4), Inches(0.35),
                            label, font_size=12, color=ACCENT, bold=True,
                            alignment=PP_ALIGN.RIGHT)
                add_textbox(slide, Inches(2.7), y, Inches(8), Inches(0.35),
                            content, font_size=12, color=TEXT_SECONDARY)
                y += Inches(0.6)
            continue

        # ── Grid (cards) slides ──
        if stype in ('grid2', 'grid3'):
            y = content_top
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=ACCENT)
                y += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                            title, font_size=28, color=TEXT_PRIMARY, bold=True)
                y += Inches(0.9)

            cards = sd.get('cards', [])
            n_cols = 3 if stype == 'grid3' else 2
            card_w = Inches(3.6) if n_cols == 3 else Inches(5.5)
            gap = Inches(0.3)
            cx = MARGIN

            for ci, card in enumerate(cards):
                card_h = Inches(2.2)
                add_shape_bg(slide, cx, y, card_w, card_h)
                add_textbox(slide, cx + Inches(0.2), y + Inches(0.2),
                            card_w - Inches(0.4), Inches(0.35),
                            card['title'], font_size=13, color=TEXT_PRIMARY,
                            bold=True)
                add_textbox(slide, cx + Inches(0.2), y + Inches(0.6),
                            card_w - Inches(0.4), Inches(1.4),
                            card['text'], font_size=10, color=TEXT_SECONDARY)
                cx += card_w + gap

                if (ci + 1) % n_cols == 0:
                    cx = MARGIN
                    y += card_h + gap
            continue

        # ── Default content slide ──
        y = content_top
        if kicker:
            add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                        kicker.upper(), font_size=11, color=ACCENT)
            y += Inches(0.4)
        if title:
            add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                        title, font_size=28, color=TEXT_PRIMARY, bold=True)
            y += Inches(0.9)
        if lead:
            add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.8),
                        lead, font_size=14, color=TEXT_SECONDARY)

    prs.save(str(output_path))
    return output_path


def main():
    if len(sys.argv) < 2:
        print("Usage: python3 convert_to_pptx.py <input.html> [output.pptx]")
        sys.exit(1)

    input_path = Path(sys.argv[1])
    if len(sys.argv) > 2:
        output_path = Path(sys.argv[2])
    else:
        output_path = input_path.with_suffix('.pptx')

    html = input_path.read_text(encoding='utf-8')
    slides_data = parse_slides(html)
    print(f"Parsed {len(slides_data)} slides")

    result = build_pptx(slides_data, output_path)
    print(f"Saved: {result}")


if __name__ == '__main__':
    main()
