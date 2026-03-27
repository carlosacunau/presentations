#!/usr/bin/env python3
"""Convert Reveal.js HTML presentations to PowerPoint (.pptx).

Usage:
    python3 convert_to_pptx-v2.py <input.html> [output.pptx]

If output is omitted, saves alongside the HTML file with same name.
Resolves image paths relative to the HTML file's directory.
Extracts theme colors from CSS variables automatically.
"""

import sys
import re
from pathlib import Path
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def hex_to_rgb(hex_str):
    """Convert '#RRGGBB' or '#RGB' to RGBColor."""
    h = hex_str.lstrip('#')
    if len(h) == 3:
        h = ''.join(c * 2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ─── Default theme (overridden by CSS vars if found) ───
DEFAULT_THEME = {
    'bg-dark': '#0A0A14',
    'bg-card': '#12121E',
    'accent': '#8B5CF6',
    'accent-light': '#A78BFA',
    'success': '#34D399',
    'warning': '#FBBF24',
    'error': '#F87171',
    'text-primary': '#E8E4F0',
    'text-secondary': '#9B95A8',
    'text-muted': '#5C5670',
    'border': '#2A2A3A',
}


class Theme:
    """Presentation theme extracted from CSS variables."""
    def __init__(self, css_vars):
        merged = {**DEFAULT_THEME, **css_vars}
        self.bg_dark = hex_to_rgb(merged['bg-dark'])
        self.bg_card = hex_to_rgb(merged['bg-card'])
        self.accent = hex_to_rgb(merged['accent'])
        self.accent_light = hex_to_rgb(merged.get('accent-light', merged['accent']))
        self.success = hex_to_rgb(merged['success'])
        self.warning = hex_to_rgb(merged['warning'])
        self.error = hex_to_rgb(merged['error'])
        self.text_primary = hex_to_rgb(merged['text-primary'])
        self.text_secondary = hex_to_rgb(merged['text-secondary'])
        self.text_muted = hex_to_rgb(merged['text-muted'])
        self.border = hex_to_rgb(merged['border'])


class HTMLTextExtractor(HTMLParser):
    """Extract plain text from HTML fragments."""
    def __init__(self):
        super().__init__()
        self._text = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in ('br',):
            self._text.append('\n')
        if tag == 'style':
            self._skip = True

    def handle_endtag(self, tag):
        if tag == 'style':
            self._skip = False

    def handle_data(self, data):
        if not self._skip:
            self._text.append(data)

    def handle_entityref(self, name):
        entities = {'amp': '&', 'lt': '<', 'gt': '>', 'nbsp': ' ',
                     'rarr': '\u2192', '#8594': '\u2192'}
        self._text.append(entities.get(name, f'&{name};'))

    def handle_charref(self, name):
        try:
            if name.startswith('x'):
                self._text.append(chr(int(name[1:], 16)))
            else:
                self._text.append(chr(int(name)))
        except ValueError:
            self._text.append(f'&#{name};')

    def get_text(self):
        return ''.join(self._text).strip()


def html_to_text(html_str):
    """Strip HTML tags, convert <br> to newlines."""
    extractor = HTMLTextExtractor()
    extractor.feed(html_str)
    return extractor.get_text()


def extract_css_vars(html):
    """Extract CSS custom properties from the HTML."""
    colors = {}
    var_pattern = re.compile(r'--([\w-]+):\s*([^;]+);')
    style_match = re.search(r'<style>(.*?)</style>', html, re.DOTALL)
    if style_match:
        for m in var_pattern.finditer(style_match.group(1)):
            val = m.group(2).strip().strip("'\"")
            if val.startswith('#'):
                colors[m.group(1)] = val
    return colors


def extract_header_logo(html):
    """Extract the header logo text from the HTML."""
    m = re.search(r'header__logo[^>]*>.*?</i>\s*(.*?)\s*</div>', html, re.DOTALL)
    if m:
        return html_to_text(m.group(1))
    m = re.search(r'header__logo[^>]*>(.*?)</div>', html, re.DOTALL)
    if m:
        return html_to_text(m.group(1))
    return ''


def parse_slides(html, html_dir):
    """Parse <section> blocks into slide data structures."""
    # Match sections with or without attributes
    sections = re.findall(r'<section([^>]*)>(.*?)</section>', html, re.DOTALL)
    slides = []

    for attrs_str, section in sections:
        slide = {'raw': section, 'elements': []}

        # Check for data-link attribute
        link_m = re.search(r'data-link="([^"]+)"', attrs_str)
        if link_m:
            slide['data_link'] = link_m.group(1)

        # Check if this is an image-only slide
        img_m = re.search(r'<img\s+src="([^"]+)"[^>]*alt="([^"]*)"', section, re.DOTALL)
        has_text_content = bool(re.search(r'<(h[12]|p\s+class="kicker"|div\s+class="(split|grid|content|metrics))', section))

        if img_m and not has_text_content:
            img_path = html_dir / img_m.group(1)
            if img_path.exists():
                slide['type'] = 'image'
                slide['image_path'] = str(img_path)
                slide['image_alt'] = img_m.group(2)
                # Check for caption below image
                caption_m = re.search(r'object-fit:\s*contain[^>]*>.*?<p[^>]*>(.*?)</p>', section, re.DOTALL)
                if caption_m:
                    slide['caption'] = html_to_text(caption_m.group(1))
                slides.append(slide)
                continue

        # Detect slide type
        if 'slide--section' in section and 'slide--center' in section:
            slide['type'] = 'section'
        elif 'arch-zone' in section:
            slide['type'] = 'arch-zones'
        elif 'split' in section and '<table' in section:
            slide['type'] = 'split-table'
        elif 'split' in section:
            slide['type'] = 'split'
        elif 'grid-3' in section:
            slide['type'] = 'grid3'
        elif 'grid-2' in section and 'card' in section:
            slide['type'] = 'grid2'
        elif 'grid-2' in section:
            slide['type'] = 'grid2-lists'
        elif 'compare' in section:
            slide['type'] = 'compare'
        elif 'pipeline' in section:
            slide['type'] = 'pipeline'
        elif 'timeline' in section:
            slide['type'] = 'timeline'
        elif 'metrics' in section:
            slide['type'] = 'metrics'
        elif 'contact-grid' in section:
            slide['type'] = 'section'
        else:
            slide['type'] = 'content'

        # Extract kicker
        kicker_m = re.search(r'class="kicker"[^>]*>(.*?)</p>', section, re.DOTALL)
        if kicker_m:
            slide['kicker'] = html_to_text(kicker_m.group(1))

        # Extract title (h1 or h2)
        title_m = re.search(r'<h[12][^>]*>(.*?)</h[12]>', section, re.DOTALL)
        if title_m:
            slide['title'] = html_to_text(title_m.group(1))

        # Extract lead text
        lead_m = re.search(r'class="lead"[^>]*>(.*?)</p>', section, re.DOTALL)
        if lead_m:
            slide['lead'] = html_to_text(lead_m.group(1))

        # Extract muted subtitle (cover slides)
        muted_m = re.search(r'class="muted"[^>]*>(.*?)</p>', section, re.DOTALL)
        if muted_m and slide['type'] == 'section':
            slide['subtitle'] = html_to_text(muted_m.group(1))

        # Extract metrics
        metrics = re.findall(
            r'class="metric__value">(.*?)</div>\s*<div class="metric__label">(.*?)</div>',
            section, re.DOTALL
        )
        if metrics:
            slide['metrics'] = [(html_to_text(v), html_to_text(l)) for v, l in metrics]
            if slide['type'] == 'content':
                slide['type'] = 'metrics'

        # Extract table
        if '<table' in section:
            rows = []
            thead = re.search(r'<thead>(.*?)</thead>', section, re.DOTALL)
            if thead:
                ths = re.findall(r'<th[^>]*>(.*?)</th>', thead.group(1), re.DOTALL)
                rows.append([html_to_text(th) for th in ths])
            tbody = re.search(r'<tbody>(.*?)</tbody>', section, re.DOTALL)
            if tbody:
                for tr in re.findall(r'<tr>(.*?)</tr>', tbody.group(1), re.DOTALL):
                    tds = re.findall(r'<td[^>]*>(.*?)</td>', tr, re.DOTALL)
                    rows.append([html_to_text(td) for td in tds])
            slide['table'] = rows

        # Extract list items
        lis = re.findall(r'<li[^>]*>(.*?)</li>', section, re.DOTALL)
        if lis:
            items = []
            for li in lis:
                strong_m = re.search(r'<strong[^>]*>(.*?)</strong>', li, re.DOTALL)
                text = html_to_text(li)
                if strong_m:
                    items.append({
                        'bold': html_to_text(strong_m.group(1)),
                        'text': text
                    })
                else:
                    items.append({'bold': '', 'text': text})
            slide['list_items'] = items

        # Extract cards (for grid2/grid3)
        if slide['type'] in ('grid2', 'grid3'):
            cards = []
            for card_m in re.finditer(r'class="card[^"]*"[^>]*>(.*?)</div>\s*</div>', section, re.DOTALL):
                card_html = card_m.group(1)
                card_title = re.search(r'class="card__title">(.*?)</h4>', card_html, re.DOTALL)
                card_texts = re.findall(r'class="card__text"[^>]*>(.*?)</p>', card_html, re.DOTALL)
                cards.append({
                    'title': html_to_text(card_title.group(1)) if card_title else '',
                    'text': ' '.join(html_to_text(t) for t in card_texts) if card_texts else ''
                })
            slide['cards'] = cards

        # Extract architecture zones
        if slide['type'] == 'arch-zones':
            zones = []
            for zone_m in re.finditer(r'class="arch-zone[^"]*"[^>]*>(.*?)</div>\s*</div>\s*</div>', section, re.DOTALL):
                zone_html = zone_m.group(0)
                label_m = re.search(r'arch-zone__label[^>]*>(.*?)</div>', zone_html, re.DOTALL)
                h4s = re.findall(r'<h4[^>]*>(.*?)</h4>', zone_html, re.DOTALL)
                ps = re.findall(r'class="card__text"[^>]*>(.*?)</p>', zone_html, re.DOTALL)
                entries = []
                for j in range(len(h4s)):
                    entries.append({
                        'title': html_to_text(h4s[j]),
                        'desc': html_to_text(ps[j]) if j < len(ps) else ''
                    })
                zones.append({
                    'label': html_to_text(label_m.group(1)) if label_m else '',
                    'entries': entries
                })
            slide['zones'] = zones

        # Extract grid-2 list columns (service catalog style)
        if slide['type'] == 'grid2-lists':
            cols = []
            for div_m in re.finditer(r'<div>\s*<span class="badge[^"]*"[^>]*>(.*?)</span>(.*?)</div>', section, re.DOTALL):
                badge = html_to_text(div_m.group(1))
                items_html = div_m.group(2)
                items = [html_to_text(li) for li in re.findall(r'<li[^>]*>(.*?)</li>', items_html, re.DOTALL)]
                cols.append({'badge': badge, 'items': items})
            slide['catalog_cols'] = cols

        # Extract terminal mockup content
        terminal_m = re.search(r'font-family:\s*var\(--font-mono\)[^>]*>(.*?)</div>\s*</div>', section, re.DOTALL)
        if terminal_m and slide['type'] == 'split':
            slide['terminal'] = html_to_text(terminal_m.group(1))

        # Extract pipeline steps
        if slide['type'] == 'pipeline':
            steps = []
            for row_m in re.finditer(r'pipeline__source">(.*?)</div>.*?pipeline__step[^"]*">(.*?)</div>', section, re.DOTALL):
                steps.append((html_to_text(row_m.group(1)), html_to_text(row_m.group(2))))
            slide['pipeline_steps'] = steps

        # Extract timeline items
        if slide['type'] == 'timeline':
            items = []
            for tm in re.finditer(r'timeline__year">(.*?)</div>.*?timeline__content">(.*?)</div>', section, re.DOTALL):
                items.append((html_to_text(tm.group(1)), html_to_text(tm.group(2))))
            slide['timeline_items'] = items

        # Extract compare columns
        if slide['type'] == 'compare':
            cols = []
            for col_m in re.finditer(r'compare__label">(.*?)</div>(.*?)</div>\s*</div>', section, re.DOTALL):
                label = html_to_text(col_m.group(1))
                items_html = col_m.group(2)
                items = [html_to_text(li) for li in re.findall(r'<li[^>]*>(.*?)</li>', items_html, re.DOTALL)]
                cols.append({'label': label, 'items': items})
            slide['compare_cols'] = cols

        # Extract callout/note boxes
        callout_m = re.search(r'border-left:\s*3px solid[^"]*">\s*<p[^>]*>(.*?)</p>', section, re.DOTALL)
        if callout_m:
            slide['callout'] = html_to_text(callout_m.group(1))

        # Extract contact grid
        contacts = re.findall(
            r'contact-item__label">(.*?)</div>\s*<div class="contact-item__value"[^>]*>(.*?)</div>',
            section, re.DOTALL
        )
        if contacts:
            slide['contacts'] = [(html_to_text(l), html_to_text(v)) for l, v in contacts]

        slides.append(slide)

    return slides


def set_slide_bg(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=None, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    """Add a text box to a slide."""
    if color is None:
        color = RGBColor(0xE8, 0xE4, 0xF0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_bg(slide, left, top, width, height, color):
    """Add a rounded rectangle background shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def build_pptx(slides_data, output_path, theme, header_logo):
    """Build the PowerPoint file from parsed slide data."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    MARGIN = Inches(0.8)
    CONTENT_W = prs.slide_width - 2 * MARGIN
    HALF_W = Inches(5.5)
    T = theme

    for i, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, T.bg_dark)

        stype = sd.get('type', 'content')
        kicker = sd.get('kicker', '')
        title = sd.get('title', '')
        lead = sd.get('lead', '')

        # ── Image-only slides ──
        if stype == 'image':
            img_path = sd['image_path']
            slide.shapes.add_picture(
                img_path, Inches(0.3), Inches(0.3),
                prs.slide_width - Inches(0.6), Inches(6.5)
            )
            if sd.get('caption'):
                add_textbox(slide, MARGIN, Inches(7.0), CONTENT_W, Inches(0.4),
                            sd['caption'], font_size=10, color=T.text_secondary,
                            alignment=PP_ALIGN.CENTER, font_name='Consolas')
            continue

        # ── Section / Cover slides ──
        if stype == 'section':
            y = Inches(2.0)
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.4),
                            kicker.upper(), font_size=11, color=T.accent,
                            alignment=PP_ALIGN.CENTER)
                y += Inches(0.5)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(1.5),
                            title, font_size=40, color=T.text_primary, bold=True,
                            alignment=PP_ALIGN.CENTER)
                y += Inches(1.5)
            if lead:
                add_textbox(slide, Inches(2.5), y, Inches(8), Inches(1.0),
                            lead, font_size=16, color=T.text_secondary,
                            alignment=PP_ALIGN.CENTER)
                y += Inches(1.2)

            # Subtitle (date line, muted text)
            if sd.get('subtitle'):
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.4),
                            sd['subtitle'], font_size=10, color=T.text_muted,
                            alignment=PP_ALIGN.CENTER)
                y += Inches(0.5)
            else:
                subtitle_m = re.search(r'font-size:\s*0\.7\d*rem[^>]*>(.*?)</p>', sd['raw'], re.DOTALL)
                if subtitle_m:
                    add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.4),
                                html_to_text(subtitle_m.group(1)), font_size=10,
                                color=T.text_muted, alignment=PP_ALIGN.CENTER)

            # Contact grid
            if sd.get('contacts'):
                y = Inches(5.2)
                cx = Inches(3.0)
                for label, value in sd['contacts']:
                    add_shape_bg(slide, cx, y, Inches(2.2), Inches(0.8), T.bg_card)
                    add_textbox(slide, cx + Inches(0.15), y + Inches(0.1),
                                Inches(1.9), Inches(0.3), label, font_size=9,
                                color=T.text_muted)
                    add_textbox(slide, cx + Inches(0.15), y + Inches(0.4),
                                Inches(1.9), Inches(0.3), value, font_size=11,
                                color=T.accent)
                    cx += Inches(2.5)
            continue

        # ── Header bar (non-section slides) ──
        header_y = Inches(0.4)
        add_textbox(slide, MARGIN, header_y, Inches(3), Inches(0.35),
                    header_logo or 'Presentation', font_size=12,
                    color=T.text_secondary, font_name='Georgia')
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN, Inches(0.85), CONTENT_W, Pt(1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = T.border
        line.line.fill.background()

        content_top = Inches(1.1)

        # ── Metrics slide ──
        if stype == 'metrics':
            y = content_top
            if kicker:
                add_textbox(slide, MARGIN, y, HALF_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=T.accent)
                y += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, y, HALF_W, Inches(1.0),
                            title, font_size=32, color=T.text_primary, bold=True)
                y += Inches(1.0)
            if lead:
                add_textbox(slide, MARGIN, y, Inches(8), Inches(0.6),
                            lead, font_size=14, color=T.text_secondary)
                y += Inches(0.8)

            if sd.get('metrics'):
                mx = MARGIN
                n_metrics = len(sd['metrics'])
                mw = min(Inches(2.8), (CONTENT_W - Inches(0.3) * (n_metrics - 1)) / n_metrics)
                for val, label in sd['metrics']:
                    add_shape_bg(slide, mx, y, mw, Inches(1.2), T.bg_card)
                    add_textbox(slide, mx + Inches(0.2), y + Inches(0.15),
                                mw - Inches(0.4), Inches(0.7), val,
                                font_size=36, color=T.text_primary, bold=True,
                                alignment=PP_ALIGN.CENTER)
                    add_textbox(slide, mx + Inches(0.2), y + Inches(0.8),
                                mw - Inches(0.4), Inches(0.3), label,
                                font_size=11, color=T.text_muted,
                                alignment=PP_ALIGN.CENTER)
                    mx += mw + Inches(0.3)

            # Cards below metrics (e.g., pricing slide)
            if 'card' in sd['raw']:
                cards = []
                for card_m in re.finditer(r'class="card[^"]*"[^>]*>(.*?)</div>\s*</div>', sd['raw'], re.DOTALL):
                    card_html = card_m.group(1)
                    ct = re.search(r'class="card__title">(.*?)</h4>', card_html, re.DOTALL)
                    cp = re.search(r'class="card__text"[^>]*>(.*?)</p>', card_html, re.DOTALL)
                    cards.append({
                        'title': html_to_text(ct.group(1)) if ct else '',
                        'text': html_to_text(cp.group(1)) if cp else ''
                    })
                if cards:
                    cy = y + Inches(1.5)
                    cx = MARGIN
                    cw = Inches(5.5)
                    for ci, card in enumerate(cards):
                        add_shape_bg(slide, cx, cy, cw, Inches(1.5), T.bg_card)
                        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.15),
                                    cw - Inches(0.4), Inches(0.3), card['title'],
                                    font_size=13, color=T.text_primary, bold=True)
                        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.5),
                                    cw - Inches(0.4), Inches(0.9), card['text'],
                                    font_size=10, color=T.text_secondary)
                        cx += cw + Inches(0.5)
            continue

        # ── Architecture zones (3-column) ──
        if stype == 'arch-zones':
            y = content_top
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=T.accent)
                y += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                            title, font_size=28, color=T.text_primary, bold=True)
                y += Inches(0.9)

            zones = sd.get('zones', [])
            zone_w = Inches(3.6)
            gap = Inches(0.3)
            zx = MARGIN

            for zone in zones:
                zone_h = Inches(4.0)
                add_shape_bg(slide, zx, y, zone_w, zone_h, T.bg_card)

                # Zone label
                add_textbox(slide, zx + Inches(0.2), y + Inches(0.15),
                            zone_w - Inches(0.4), Inches(0.3),
                            zone['label'].upper(), font_size=9,
                            color=T.accent, bold=True)

                zy = y + Inches(0.55)
                for entry in zone['entries']:
                    add_textbox(slide, zx + Inches(0.2), zy,
                                zone_w - Inches(0.4), Inches(0.3),
                                entry['title'], font_size=12,
                                color=T.text_primary, bold=True)
                    zy += Inches(0.35)
                    if entry['desc']:
                        add_textbox(slide, zx + Inches(0.2), zy,
                                    zone_w - Inches(0.4), Inches(0.5),
                                    entry['desc'], font_size=10,
                                    color=T.text_secondary)
                        zy += Inches(0.5)

                zx += zone_w + gap

            # Callout below zones
            if sd.get('callout'):
                cy = y + Inches(4.2)
                add_shape_bg(slide, MARGIN, cy, CONTENT_W, Inches(0.6), T.bg_card)
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, MARGIN, cy, Pt(3), Inches(0.6)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = T.accent
                bar.line.fill.background()
                add_textbox(slide, MARGIN + Inches(0.2), cy + Inches(0.1),
                            CONTENT_W - Inches(0.4), Inches(0.4),
                            sd['callout'], font_size=10, color=T.text_secondary)
            continue

        # ── Grid-2 lists (service catalog) ──
        if stype == 'grid2-lists':
            y = content_top
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=T.accent)
                y += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                            title, font_size=28, color=T.text_primary, bold=True)
                y += Inches(0.8)

            cols = sd.get('catalog_cols', [])
            if not cols and sd.get('list_items'):
                # Fallback: split list items into two columns
                items = sd['list_items']
                mid = len(items) // 2
                cols = [
                    {'badge': 'Column 1', 'items': [i['text'] for i in items[:mid]]},
                    {'badge': 'Column 2', 'items': [i['text'] for i in items[mid:]]}
                ]

            cx = MARGIN
            col_w = Inches(5.5)
            for ci, col in enumerate(cols):
                add_shape_bg(slide, cx, y, col_w, Inches(4.5), T.bg_card)
                badge_color = T.success if ci == 0 else T.accent
                add_textbox(slide, cx + Inches(0.2), y + Inches(0.15),
                            col_w - Inches(0.4), Inches(0.3),
                            col['badge'], font_size=10, color=badge_color, bold=True)
                iy = y + Inches(0.55)
                for item in col['items']:
                    add_textbox(slide, cx + Inches(0.3), iy,
                                col_w - Inches(0.5), Inches(0.3),
                                '\u2713  ' + item if ci == 0 else '\u2022  ' + item,
                                font_size=11, color=T.text_secondary)
                    iy += Inches(0.38)
                cx += col_w + Inches(0.5)
            continue

        # ── Split layouts ──
        if stype in ('split', 'split-table'):
            left_x = MARGIN
            right_x = Inches(6.8)
            y = content_top

            # Left column
            if kicker:
                add_textbox(slide, left_x, y, HALF_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=T.accent)
                y += Inches(0.4)
            if title:
                add_textbox(slide, left_x, y, HALF_W, Inches(1.2),
                            title, font_size=28, color=T.text_primary, bold=True)
                y += Inches(1.1)
            if lead:
                add_textbox(slide, left_x, y, HALF_W, Inches(1.5),
                            lead, font_size=13, color=T.text_secondary)
                y += Inches(1.0)

            # Left-side list (for splits that have list on left, like Morning Routine)
            left_list_m = re.search(r'</p>\s*<ul class="list"[^>]*>(.*?)</ul>', sd['raw'], re.DOTALL)
            if left_list_m and sd.get('terminal'):
                left_items = re.findall(r'<li[^>]*>(.*?)</li>', left_list_m.group(1), re.DOTALL)
                for li_text in left_items:
                    add_textbox(slide, left_x + Inches(0.2), y,
                                HALF_W - Inches(0.4), Inches(0.3),
                                '\u2022 ' + html_to_text(li_text),
                                font_size=11, color=T.text_secondary)
                    y += Inches(0.35)

            # Callout box on left
            if sd.get('callout'):
                add_shape_bg(slide, left_x, y, HALF_W, Inches(0.8), T.bg_card)
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left_x, y, Pt(3), Inches(0.8)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = T.accent
                bar.line.fill.background()
                add_textbox(slide, left_x + Inches(0.2), y + Inches(0.1),
                            HALF_W - Inches(0.4), Inches(0.6),
                            sd['callout'], font_size=10, color=T.text_secondary)

            # Right column
            ry = content_top + Inches(0.2)

            # Terminal mockup
            if sd.get('terminal'):
                term_text = sd['terminal']
                add_shape_bg(slide, right_x, ry, Inches(5.5), Inches(5.0),
                             RGBColor(0x08, 0x08, 0x0F))
                add_textbox(slide, right_x + Inches(0.2), ry + Inches(0.2),
                            Inches(5.1), Inches(4.6), term_text,
                            font_size=9, color=T.text_secondary,
                            font_name='Consolas')
            elif sd.get('table'):
                table_data = sd['table']
                rows = len(table_data)
                cols = len(table_data[0]) if table_data else 2
                tw = Inches(5.5)
                th = Inches(0.4) * rows
                tbl = slide.shapes.add_table(rows, cols, right_x, ry, tw, th).table

                for r, row in enumerate(table_data):
                    for c, cell_text in enumerate(row):
                        cell = tbl.cell(r, c)
                        cell.text = cell_text
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A) if r == 0 else T.bg_card
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(10)
                            paragraph.font.color.rgb = T.text_primary if r == 0 else T.text_secondary
                            paragraph.font.bold = (r == 0)
                            paragraph.font.name = 'Calibri'

            elif sd.get('list_items'):
                for item in sd['list_items']:
                    text = item['text']
                    txBox = slide.shapes.add_textbox(right_x, ry, Inches(5.5), Inches(0.8))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]

                    if item['bold']:
                        run_b = p.add_run()
                        run_b.text = item['bold'] + '\n'
                        run_b.font.size = Pt(12)
                        run_b.font.color.rgb = T.text_primary
                        run_b.font.bold = True
                        run_b.font.name = 'Calibri'
                        rest = text.replace(item['bold'], '', 1).strip()
                        if rest:
                            run_r = p.add_run()
                            run_r.text = rest
                            run_r.font.size = Pt(10)
                            run_r.font.color.rgb = T.text_muted
                            run_r.font.name = 'Calibri'
                    else:
                        p.text = text
                        p.font.size = Pt(11)
                        p.font.color.rgb = T.text_secondary
                        p.font.name = 'Calibri'
                    ry += Inches(0.85)
            continue

        # ── Compare slide ──
        if stype == 'compare':
            y = content_top
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=T.accent)
                y += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                            title, font_size=28, color=T.text_primary, bold=True)
                y += Inches(0.9)

            cols = sd.get('compare_cols', [])
            cx = MARGIN
            col_w = Inches(5.5)
            for ci, col in enumerate(cols):
                is_after = ci == 1
                add_shape_bg(slide, cx, y, col_w, Inches(4.0),
                             RGBColor(0x0F, 0x14, 0x0F) if is_after else T.bg_card)
                label_color = T.success if is_after else T.warning
                add_textbox(slide, cx + Inches(0.2), y + Inches(0.15),
                            col_w - Inches(0.4), Inches(0.35),
                            col['label'], font_size=13, color=label_color, bold=True)
                iy = y + Inches(0.6)
                for item in col['items']:
                    bullet_color = T.success if is_after else T.text_secondary
                    add_textbox(slide, cx + Inches(0.3), iy,
                                col_w - Inches(0.6), Inches(0.35),
                                '\u2022 ' + item, font_size=11, color=bullet_color)
                    iy += Inches(0.45)
                cx += col_w + Inches(0.5)
            continue

        # ── Pipeline slide ──
        if stype == 'pipeline':
            y = content_top
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=T.accent)
                y += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                            title, font_size=28, color=T.text_primary, bold=True)
                y += Inches(0.9)

            for source, step in sd.get('pipeline_steps', []):
                add_shape_bg(slide, MARGIN, y, CONTENT_W, Inches(0.55), T.bg_card)
                add_textbox(slide, MARGIN + Inches(0.2), y + Inches(0.08),
                            Inches(2.5), Inches(0.35), source,
                            font_size=10, color=T.accent, bold=True)
                add_textbox(slide, Inches(3.8), y + Inches(0.08),
                            Inches(0.4), Inches(0.35), '\u2192',
                            font_size=12, color=T.text_muted,
                            alignment=PP_ALIGN.CENTER)
                add_textbox(slide, Inches(4.3), y + Inches(0.08),
                            Inches(8), Inches(0.35), step,
                            font_size=10, color=T.text_secondary)
                y += Inches(0.65)
            continue

        # ── Timeline slide ──
        if stype == 'timeline':
            y = content_top
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=T.accent)
                y += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                            title, font_size=28, color=T.text_primary, bold=True)
                y += Inches(1.0)

            for idx, (label, content) in enumerate(sd.get('timeline_items', [])):
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, Inches(2.2), y + Inches(0.12), Pt(10), Pt(10)
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = T.accent
                dot.line.fill.background()
                if idx < len(sd['timeline_items']) - 1:
                    vline = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(2.24), y + Inches(0.35), Pt(2), Inches(0.4)
                    )
                    vline.fill.solid()
                    vline.fill.fore_color.rgb = T.border
                    vline.line.fill.background()

                add_textbox(slide, MARGIN, y, Inches(1.4), Inches(0.35),
                            label, font_size=12, color=T.accent, bold=True,
                            alignment=PP_ALIGN.RIGHT)
                add_textbox(slide, Inches(2.7), y, Inches(8), Inches(0.35),
                            content, font_size=12, color=T.text_secondary)
                y += Inches(0.6)
            continue

        # ── Grid (cards) slides ──
        if stype in ('grid2', 'grid3'):
            y = content_top
            if kicker:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                            kicker.upper(), font_size=11, color=T.accent)
                y += Inches(0.4)
            if title:
                add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                            title, font_size=28, color=T.text_primary, bold=True)
                y += Inches(0.9)

            cards = sd.get('cards', [])
            n_cols = 3 if stype == 'grid3' else 2
            card_w = Inches(3.6) if n_cols == 3 else Inches(5.5)
            gap = Inches(0.3)
            cx = MARGIN

            for ci, card in enumerate(cards):
                card_h = Inches(2.2)
                add_shape_bg(slide, cx, y, card_w, card_h, T.bg_card)
                add_textbox(slide, cx + Inches(0.2), y + Inches(0.2),
                            card_w - Inches(0.4), Inches(0.35),
                            card['title'], font_size=13, color=T.text_primary,
                            bold=True)
                add_textbox(slide, cx + Inches(0.2), y + Inches(0.6),
                            card_w - Inches(0.4), Inches(1.4),
                            card['text'], font_size=10, color=T.text_secondary)
                cx += card_w + gap

                if (ci + 1) % n_cols == 0:
                    cx = MARGIN
                    y += card_h + gap
            continue

        # ── Default content slide ──
        y = content_top
        if kicker:
            add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                        kicker.upper(), font_size=11, color=T.accent)
            y += Inches(0.4)
        if title:
            add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                        title, font_size=28, color=T.text_primary, bold=True)
            y += Inches(0.9)
        if lead:
            add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.8),
                        lead, font_size=14, color=T.text_secondary)

    prs.save(str(output_path))
    return output_path


def main():
    if len(sys.argv) < 2:
        print("Usage: python3 convert_to_pptx-v2.py <input.html> [output.pptx]")
        sys.exit(1)

    input_path = Path(sys.argv[1])
    if len(sys.argv) > 2:
        output_path = Path(sys.argv[2])
    else:
        output_path = input_path.with_suffix('.pptx')

    html = input_path.read_text(encoding='utf-8')
    html_dir = input_path.parent

    # Extract theme from CSS vars
    css_vars = extract_css_vars(html)
    theme = Theme(css_vars)

    # Extract header logo text
    header_logo = extract_header_logo(html)

    slides_data = parse_slides(html, html_dir)
    print(f"Parsed {len(slides_data)} slides (theme: {len(css_vars)} CSS vars, header: '{header_logo}')")

    result = build_pptx(slides_data, output_path, theme, header_logo)
    print(f"Saved: {result}")


if __name__ == '__main__':
    main()
